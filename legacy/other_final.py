from pathlib import Path
from datetime import datetime
import csv
import re
import sys
import html
import subprocess
import tempfile

from openpyxl import Workbook
from openpyxl.styles import Font, Alignment


# ============================================================
# CONFIGURATION
# ============================================================

ROOT = Path(sys.argv[1]) if len(sys.argv) > 1 else Path.cwd()

OUTPUT = ROOT / "Chaput_Other_Documents.xlsx"

EXCLUDED_DIRS = {
    "_doc_catalog",
    "raw_text",
}


# Files already represented by the Word catalog
EXTENSIONS = {
    ".pdf",
    ".txt",
    ".htm",
    ".html",
    ".rtf",
    ".odt",
    ".ppt",
    ".pptx",
}


# ============================================================
# TEXT EXTRACTION
# ============================================================

def clean_text(text):
    """Collapse whitespace and remove characters Excel cannot store."""

    # Remove XML/Excel-illegal control characters
    text = re.sub(
        r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]",
        "",
        text,
    )

    text = re.sub(r"\s+", " ", text)

    return text.strip()


def extract_pdf(path):
    try:
        from pypdf import PdfReader

        reader = PdfReader(path)
        parts = []

        for page in reader.pages:
            text = page.extract_text()

            if text:
                parts.append(text)

            if len(" ".join(parts)) >= 500:
                break

        return clean_text(" ".join(parts))[:500]

    except Exception:
        return ""


def extract_txt(path):
    try:
        text = path.read_text(
            encoding="utf-8",
            errors="ignore",
        )

        return clean_text(text)[:500]

    except Exception:
        return ""


def extract_rtf(path):
    try:
        text = path.read_text(
            encoding="utf-8",
            errors="ignore",
        )

        # Basic RTF cleanup
        text = re.sub(r"\\'[0-9a-fA-F]{2}", "", text)
        text = re.sub(r"\\[a-zA-Z]+\d* ?", "", text)
        text = text.replace("{", "").replace("}", "")

        return clean_text(text)[:500]

    except Exception:
        return ""


def extract_html(path):
    try:
        from bs4 import BeautifulSoup

        raw = path.read_text(
            encoding="utf-8",
            errors="ignore",
        )

        soup = BeautifulSoup(raw, "html.parser")

        title = soup.title.get_text(" ", strip=True) if soup.title else ""

        text = soup.get_text(" ", strip=True)
        text = clean_text(html.unescape(text))

        if title:
            return clean_text(
                f"{title} — {text}"
            )[:500]

        return text[:500]

    except Exception:
        return ""


def extract_ppt(path):
    try:
        from pptx import Presentation

        presentation = Presentation(path)

        parts = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text = shape.text.strip()

                    if text:
                        parts.append(text)

            if len(" ".join(parts)) >= 500:
                break

        return clean_text(" ".join(parts))[:500]

    except Exception:
        return ""


def extract_odt(path):
    """
    ODT is a ZIP/XML document.
    """

    try:
        import zipfile
        from xml.etree import ElementTree as ET

        with zipfile.ZipFile(path) as z:

            xml = z.read("content.xml")

        root = ET.fromstring(xml)

        parts = []

        for element in root.iter():

            if element.text:
                parts.append(element.text)

        return clean_text(" ".join(parts))[:500]

    except Exception:
        return ""


def extract_preview(path):

    ext = path.suffix.lower()

    if ext == ".pdf":
        return extract_pdf(path)

    if ext == ".txt":
        return extract_txt(path)

    if ext == ".rtf":
        return extract_rtf(path)

    if ext in {".htm", ".html"}:
        return extract_html(path)

    if ext in {".ppt", ".pptx"}:
        return extract_ppt(path)

    if ext == ".odt":
        return extract_odt(path)

    return ""


# ============================================================
# FIND FILES
# ============================================================

files = []

for path in ROOT.rglob("*"):

    if not path.is_file():
        continue

    if any(part in EXCLUDED_DIRS for part in path.parts):
        continue

    if path.suffix.lower() not in EXTENSIONS:
        continue

    files.append(path)


print()
print("=" * 60)
print("CHAPUT OTHER DOCUMENT CATALOG")
print("=" * 60)
print()
print(f"Found {len(files):,} files")
print()


# ============================================================
# BUILD RECORDS
# ============================================================

records = []

for i, path in enumerate(files, 1):

    relative = path.relative_to(ROOT)

    file_location = str(relative).replace("\\", "/")

    try:

        stat = path.stat()

        modified = datetime.fromtimestamp(
            stat.st_mtime
        ).isoformat(
            sep=" ",
            timespec="seconds",
        )

    except Exception:

        modified = ""

    preview = extract_preview(path)

    records.append({
        "type": path.suffix.lower().lstrip(".").upper(),
        "filename": path.name,
        "file_location": file_location,
        "modified": modified,
        "preview": preview,
    })

    if i % 25 == 0 or i == len(files):

        print(
            f"Processed {i:,}/{len(files):,}"
        )


# ============================================================
# SORT
# ============================================================

records.sort(
    key=lambda r: (
        r["type"].lower(),
        r["filename"].lower(),
    )
)


# ============================================================
# CREATE WORKBOOK
# ============================================================

wb = Workbook()

ws = wb.active
ws.title = "Other Documents"


headers = [
    "Type",
    "Filename",
    "Location",
    "Modified",
    "Preview",
]

ws.append(headers)


# ============================================================
# ADD DOCUMENTS
# ============================================================

for record in records:

    ws.append([
        record["type"],
        record["filename"],
        record["file_location"],
        record["modified"],
        record["preview"],
    ])

    row = ws.max_row

    # --------------------------------------------------------
    # Filename → original file
    # --------------------------------------------------------

    filename_cell = ws.cell(row, 2)

    filename_cell.hyperlink = (
        f"./{record['file_location']}"
    )

    filename_cell.font = Font(
        color="0563C1",
        underline="single",
    )

    # --------------------------------------------------------
    # File Location → original file
    # --------------------------------------------------------

    location_cell = ws.cell(row, 3)

    location_cell.hyperlink = (
        f"./{record['file_location']}"
    )

    location_cell.font = Font(
        color="0563C1",
        underline="single",
    )

    # --------------------------------------------------------
    # Preview
    # --------------------------------------------------------

    ws.cell(row, 5).alignment = Alignment(
        vertical="top",
        wrap_text=True,
    )


# ============================================================
# FORMAT
# ============================================================

# Freeze header
ws.freeze_panes = "A2"

# Filter
ws.auto_filter.ref = ws.dimensions


# Header
for cell in ws[1]:

    cell.font = Font(bold=True)

    cell.alignment = Alignment(
        vertical="center"
    )


# Column widths
ws.column_dimensions["A"].width = 8
ws.column_dimensions["B"].width = 45
ws.column_dimensions["C"].width = 10
ws.column_dimensions["D"].width = 22
ws.column_dimensions["E"].width = 100


ws.row_dimensions[1].height = 24


# ============================================================
# SAVE
# ============================================================

wb.save(OUTPUT)


print()
print("=" * 60)
print("CATALOG CREATED")
print("=" * 60)
print()
print(f"Documents: {len(records):,}")
print(f"Output:    {OUTPUT}")
print()
print("Sorted by: Type → Filename")
print("Clickable: Filename, File Location")
print("Header row: frozen")
print("Preview: extracted where possible")